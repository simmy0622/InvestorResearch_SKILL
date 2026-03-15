from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG      = RGBColor(0x0F, 0x11, 0x17)
CARD    = RGBColor(0x1A, 0x1D, 0x27)
BLUE    = RGBColor(0x4E, 0x9A, 0xF5)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
ORANGE  = RGBColor(0xF5, 0xA6, 0x23)
RED     = RGBColor(0xE8, 0x5D, 0x75)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x8A, 0x90, 0xA0)
LGRAY   = RGBColor(0xCB, 0xD0, 0xDC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(s):
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG

def rect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    if sh.adjustments: sh.adjustments[0] = 0.05
    return sh

def txt(s, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT, fn="Helvetica Neue"):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = fn; p.alignment = a
    return tb

def dot(s, l, t, sz, c, text="", tsz=20):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    p = sh.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(tsz); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════
# P1 — 封面
# ═══════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s, Inches(0.7), Inches(2.9), Inches(0.07), Inches(1.4), BLUE)
txt(s, Inches(1.2), Inches(2.7), Inches(8), Inches(1), "TradeTriad", sz=54, c=WHITE, b=True)
txt(s, Inches(1.2), Inches(3.7), Inches(8), Inches(0.6), "A 股波段交易情报管线", sz=24, c=LGRAY)
txt(s, Inches(1.2), Inches(5.5), Inches(6), Inches(0.4),
    "v3.0  ·  OpenClaw  ·  5 个 Skill  ·  全自动", sz=13, c=GRAY)

cs = [BLUE, GREEN, ORANGE, RED, PURPLE]
ls = ["抓取", "分析", "计划", "风控", "展示"]
for i in range(5):
    x = Inches(8.2 + i * 1.0)
    dot(s, x, Inches(3.2), Inches(0.6), cs[i], str(i+1), 20)
    txt(s, x - Inches(0.08), Inches(3.95), Inches(0.76), Inches(0.35), ls[i], sz=12, c=GRAY, a=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, x + Inches(0.6), Inches(3.25), Inches(0.4), Inches(0.5), "›", sz=22, c=GRAY, a=PP_ALIGN.CENTER)

# ═══════════════════════════════════
# P2 — 五阶段管线
# ═══════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6), "五阶段管线", sz=32, c=WHITE, b=True)
rect(s, Inches(0.8), Inches(0.95), Inches(1.6), Inches(0.04), BLUE)

stages = [
    ("1", "抓取", "三源抓取 → 去重 → 六维评分", BLUE),
    ("2", "分析", "宏观×行业×资产 三层交叉验证", GREEN),
    ("3", "计划", "标的映射 + 入场/止损/失效条件", ORANGE),
    ("4", "风控", "7 条 veto 规则 独立审查", RED),
    ("5", "展示", "自动打包 → 浏览器 Artifact Desk", PURPLE),
]

cw = Inches(2.28); ch = Inches(3.2); gap = Inches(0.22); sx = Inches(0.5); sy = Inches(1.6)
for i, (n, name, desc, c) in enumerate(stages):
    x = sx + i * (cw + gap)
    rect(s, x, sy, cw, ch, CARD)
    dot(s, x + Inches(0.78), sy + Inches(0.35), Inches(0.6), c, n, 22)
    txt(s, x, sy + Inches(1.15), cw, Inches(0.5), name, sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.15), sy + Inches(1.8), cw - Inches(0.3), Inches(1.2), desc, sz=13, c=GRAY, a=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, x + cw, sy + Inches(1.2), gap, Inches(0.5), "→", sz=20, c=GRAY, a=PP_ALIGN.CENTER)

rect(s, Inches(0.5), Inches(5.2), Inches(12.33), Inches(0.6), CARD)
txt(s, Inches(0.8), Inches(5.25), Inches(11.8), Inches(0.5),
    "triad-orchestrator 统一调度  ·  一个 runId 贯穿全程  ·  阶段失败自动记录  ·  UI 强制刷新",
    sz=13, c=GRAY, a=PP_ALIGN.CENTER)

# ═══════════════════════════════════
# P3 — 评分 + 风控（左右分栏）
# ═══════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)

# Left: scoring
txt(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6), "六维评分", sz=32, c=WHITE, b=True)
rect(s, Inches(0.8), Inches(0.95), Inches(1.4), Inches(0.04), BLUE)

dims = [
    ("时效性", "×2.0", BLUE),   ("影响力", "×2.5", BLUE),
    ("确定性", "×2.0", GREEN),   ("相关度", "×1.5", GREEN),
    ("意外度", "×1.5", ORANGE),  ("可交易", "×1.5", ORANGE),
]
for i, (nm, wt, c) in enumerate(dims):
    y = Inches(1.35 + i * 0.62)
    rect(s, Inches(0.8), y, Inches(5.3), Inches(0.5), CARD)
    rect(s, Inches(0.95), y + Inches(0.16), Inches(0.18), Inches(0.18), c)
    txt(s, Inches(1.3), y + Inches(0.06), Inches(2.5), Inches(0.38), nm, sz=15, c=LGRAY)
    txt(s, Inches(4.8), y + Inches(0.06), Inches(1.1), Inches(0.38), wt, sz=15, c=c, b=True, a=PP_ALIGN.RIGHT)

rect(s, Inches(0.8), Inches(5.15), Inches(5.3), Inches(0.5), CARD)
txt(s, Inches(1.0), Inches(5.2), Inches(4.9), Inches(0.4),
    "S≥85  A≥68 → 进入分析    B≥45  C<45 → 归档", sz=13, c=GRAY, a=PP_ALIGN.CENTER)

# Right: veto
txt(s, Inches(7.0), Inches(0.4), Inches(5), Inches(0.6), "七条风控", sz=32, c=WHITE, b=True)
rect(s, Inches(7.0), Inches(0.95), Inches(1.4), Inches(0.04), RED)

vetos = [
    ("V1", "缺少标的",      "watch"),
    ("V2", "缺少失效条件",   "watch"),
    ("V3", "缺少技术确认",   "watch"),
    ("V4", "来源冲突",      "watch/rej"),
    ("V5", "流动性风险",    "rej/watch"),
    ("V6", "同向敞口集中",   "watch"),
    ("V7", "不可证伪",      "rejected"),
]
for i, (cd, nm, rs) in enumerate(vetos):
    y = Inches(1.35 + i * 0.62)
    rect(s, Inches(7.0), y, Inches(5.5), Inches(0.5), CARD)
    txt(s, Inches(7.2), y + Inches(0.06), Inches(0.5), Inches(0.38), cd, sz=13, c=RED, b=True, fn="Menlo")
    txt(s, Inches(7.8), y + Inches(0.06), Inches(2.8), Inches(0.38), nm, sz=15, c=LGRAY)
    txt(s, Inches(10.8), y + Inches(0.06), Inches(1.5), Inches(0.38), rs, sz=12, c=GRAY, a=PP_ALIGN.RIGHT, fn="Menlo")

rect(s, Inches(7.0), Inches(5.15), Inches(5.5), Inches(0.5), CARD)
dcs = [(GREEN, "approved"), (ORANGE, "watch"), (RED, "rejected")]
for j, (dc, dl) in enumerate(dcs):
    dx = Inches(7.5 + j * 1.8)
    rect(s, dx, Inches(5.27), Inches(0.16), Inches(0.16), dc)
    txt(s, dx + Inches(0.25), Inches(5.22), Inches(1.3), Inches(0.3), dl, sz=12, c=GRAY)

# ═══════════════════════════════════
# P4 — 结尾
# ═══════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1), "TradeTriad", sz=52, c=WHITE, b=True, a=PP_ALIGN.CENTER)
rect(s, Inches(6.0), Inches(3.6), Inches(1.3), Inches(0.04), BLUE)
txt(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.7),
    "从信息到判断，从判断到计划，从计划到风控", sz=22, c=GRAY, a=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.4),
    "5 Skills  ·  6 维评分  ·  7 条风控  ·  全自动管线", sz=14, c=GRAY, a=PP_ALIGN.CENTER)

# ── Save ──
out = "/Users/shinmuchen/Desktop/triad_skills/TradeTriad_Intro.pptx"
prs.save(out)
print(f"Done: {out}")
